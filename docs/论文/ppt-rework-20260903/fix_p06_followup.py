# -*- coding: utf-8 -*-
"""第6页三处收尾修正：恢复来源行、缩短批次清单、HumanEval 类目全显。"""
import copy, sys, io, zipfile, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import lxml.etree as etree

SRC = "deliver/基于大语言模型的代码安全分析系统_30页修订版_20260903.pptx"
INK2 = RGBColor(0x5B, 0x67, 0x70)
INK3 = RGBColor(0x94, 0x9E, 0xA6)

def rewrite_shape_text(sh, lines, size, color):
    tf = sh.text_frame
    proto = None
    for p0 in tf.paragraphs:
        for r0 in p0.runs:
            proto = copy.deepcopy(r0._r); break
        if proto is not None: break
    txBody = tf._txBody
    for p in txBody.findall(qn('a:p'))[1:]:
        txBody.remove(p)
    p1 = txBody.find(qn('a:p'))
    for child in list(p1):
        if child.tag != qn('a:pPr'):
            p1.remove(child)
    paras = [p1]
    for _ in range(len(lines) - 1):
        np = copy.deepcopy(p1)
        for child in list(np):
            if child.tag != qn('a:pPr'):
                np.remove(child)
        txBody.append(np); paras.append(np)
    for pel, line in zip(paras, lines):
        r = copy.deepcopy(proto)
        for t in r.findall(qn('a:t')):
            r.remove(t)
        t = etree.SubElement(r, qn('a:t')); t.text = line
        rPr = r.find(qn('a:rPr'))
        if rPr is None:
            rPr = etree.Element(qn('a:rPr'), {'lang': 'zh-CN'}); r.insert(0, rPr)
        rPr.set('sz', str(int(size * 100))); rPr.set('b', '0')
        for e in rPr.findall(qn('a:solidFill')):
            rPr.remove(e)
        sf = etree.SubElement(rPr, qn('a:solidFill'))
        clr = etree.SubElement(sf, qn('a:srgbClr')); clr.set('val', str(color))
        pel.append(r)

prs = Presentation(SRC)
s6 = prs.slides[5]
for sh in s6.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text
    if t.startswith("2026.02 统一框架批次五模型"):
        # 误改的底部来源行 -> 恢复为来源注释
        rewrite_shape_text(sh, [
            "来源：SWE-bench Verified 各厂官方发布（2026-09 检索）；统一框架为 mini-SWE-agent v2.0.0（swebench.com，2026-02-17），自报框架普遍高 4\u20136 分，标 * 分列",
            "官方分：GPT-4o 33.2 / Claude 3.5 Sonnet 33.4\u219249.0 / DeepSeek-V3 42.0 / Gemini 2.5 Pro 63.8 / Claude Opus 4 72.5 / Qwen3-Coder 69.6 / GPT-5 74.9 / GLM-4.6 68.0 / Claude Opus 4.5 80.9",
            "HumanEval 取各代官方报告（Codex 28.8 \u2192 Qwen2.5-Coder-32B 92.7）；Grok 4 为 xAI 自报口径（72\u201375）",
        ], 8.5, INK3)
        print("底部来源行已恢复")
    elif t.startswith("2026.02 统一框架批次（mini-SWE-agent"):
        rewrite_shape_text(sh,
            ["2026.02 统一框架批次：① Claude Opus 4.5 76.8  ② Gemini 3 Flash 75.8  ③ MiniMax M2.5 75.8  ④ Claude Opus 4.6 75.6  ⑤ GPT-5.2 72.8"],
            7, INK2)
        print("批次清单已缩短")

prs.save(SRC)

# HumanEval 图：类目标签 7pt + 强制逐个显示
zin = zipfile.ZipFile(SRC)
target = None
for n in zin.namelist():
    if 'charts/chart' in n and n.endswith('.xml'):
        x = zin.read(n).decode('utf-8')
        if 'Qwen2.5-Coder' in x:
            target = n; break
zin.close()
if target:
    tmp = SRC + ".tmp"
    zin = zipfile.ZipFile(SRC)
    zout = zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED)
    for item in zin.infolist():
        data = zin.read(item.filename)
        if item.filename == target:
            x = data.decode('utf-8')
            # 类目轴 txPr 字号 800 -> 700（仅 catAx 段）
            i0 = x.find('<c:catAx>')
            i1 = x.find('</c:catAx>') + len('</c:catAx>')
            seg = x[i0:i1]
            seg2 = seg.replace('sz="800"', 'sz="700"')
            if '<c:tickLblSkip' not in seg2:
                seg2 = seg2.replace('</c:catAx>', '<c:tickLblSkip val="1"/></c:catAx>')
            x = x[:i0] + seg2 + x[i1:]
            data = x.encode('utf-8')
        zout.writestr(item, data)
    zout.close(); zin.close()
    os.replace(tmp, SRC)
    print("HumanEval 类目轴已修正:", target)
print("DONE")
