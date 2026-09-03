# -*- coding: utf-8 -*-
"""提取用户改过的 PPTX 每页文本，与在线稿文本对比，输出差异。"""
import re, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
import xml.etree.ElementTree as ET

# 1. 用户 PPTX
prs = Presentation(r"export/基于大语言模型的代码安全分析系统.pptx")
user_pages = []
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t: texts.append(t)
    user_pages.append(texts)

# 2. 在线稿 live_now.xml
ns = {"s": "https://www.larkoffice.com/sml/2.0"}
tree = ET.parse("live_now.xml")
slides = tree.getroot().findall(".//s:slide", ns)
live_pages = []
for sl in slides:
    texts = []
    for c in sl.findall(".//s:p", ns):
        t = "".join(t.text or "" for t in c.iter() if t.tag.endswith("}t") or t.tag=="t")
        # 简单拼接
        t = "".join(c.itertext()).strip()
        if t: texts.append(t)
    live_pages.append(texts)

print("user pages:", len(user_pages), " live pages:", len(live_pages))
for i in range(max(len(user_pages), len(live_pages))):
    u = user_pages[i] if i < len(user_pages) else []
    l = live_pages[i] if i < len(live_pages) else []
    us = "\n".join(u); ls = "\n".join(l)
    if us != ls:
        print(f"\n===== 第 {i+1} 页有差异 =====")
        print("--- 用户版 ---")
        print(us[:1400])
