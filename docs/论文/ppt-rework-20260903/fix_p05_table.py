# -*- coding: utf-8 -*-
"""第5页改造：图表 → 四方法 × 六口径统一对比表（素材库 2026-09-02 官方口径）。
数据源：素材库_论文写作素材收集.md 1.1 层1/层2 + unified_score_table_20260903.md + exp_07 Semgrep 提取。
"""
import sys, io, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import lxml.etree as etree

SRC = "deliver/基于大语言模型的代码安全分析系统_30页修订版_20260903.pptx"

INK   = RGBColor(0x2B, 0x2F, 0x36)
INK2  = RGBColor(0x5B, 0x67, 0x70)
INK3  = RGBColor(0x94, 0x9E, 0xA6)
BLUE  = RGBColor(0x4A, 0x7F, 0xA5)
BLUE_D= RGBColor(0x34, 0x5E, 0x80)
BLUE_L= RGBColor(0xE8, 0xF0, 0xF6)
GOLD  = RGBColor(0xD9, 0x96, 0x2E)
RED   = RGBColor(0xC0, 0x50, 0x4D)
TEAL  = RGBColor(0x6F, 0xA3, 0x9B)
PAPER = RGBColor(0xFA, 0xF8, 0xF3)
MIST  = RGBColor(0xF4, 0xF7, 0xFA)
LINEC = RGBColor(0xE2, 0xE6, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "思源黑体"
MONO  = "Roboto Mono"

def rewrite(sh, lines, size, color, bold=None):
    tf = sh.text_frame
    proto = None
    for p0 in tf.paragraphs:
        for r0 in p0.runs:
            proto = copy.deepcopy(r0._r); break
        if proto is not None: break
    txBody = tf._txBody
    for p in txBody.findall(qn('a:p'))[1:]: txBody.remove(p)
    p1 = txBody.find(qn('a:p'))
    for child in list(p1):
        if child.tag != qn('a:pPr'): p1.remove(child)
    paras = [p1]
    for _ in range(len(lines) - 1):
        np = copy.deepcopy(p1)
        for child in list(np):
            if child.tag != qn('a:pPr'): np.remove(child)
        txBody.append(np); paras.append(np)
    for pel, line in zip(paras, lines):
        r = copy.deepcopy(proto)
        for t in r.findall(qn('a:t')): r.remove(t)
        t = etree.SubElement(r, qn('a:t')); t.text = line
        rPr = r.find(qn('a:rPr'))
        if rPr is None:
            rPr = etree.Element(qn('a:rPr'), {'lang': 'zh-CN'}); r.insert(0, rPr)
        rPr.set('sz', str(int(size * 100)))
        if bold is not None: rPr.set('b', '1' if bold else '0')
        for e in rPr.findall(qn('a:solidFill')): rPr.remove(e)
        sf = etree.SubElement(rPr, qn('a:solidFill'))
        clr = etree.SubElement(sf, qn('a:srgbClr')); clr.set('val', str(color))
        pel.append(r)

prs = Presentation(SRC)
s5 = prs.slides[4]

# ---- 1) 删除旧图表与零值注释 ----
for sh in list(s5.shapes):
    if sh.has_chart:
        sh._element.getparent().remove(sh._element)
        print("旧图表已删除")
    elif sh.has_text_frame and sh.text_frame.text.strip().startswith("注：误报率"):
        sh._element.getparent().remove(sh._element)
        print("零值注释已删除")

# ---- 2) 插入对比表 ----
HEAD = ["方法", "召回率", "误报率", "准确率", "strict 归因", "CVE-20 召回", "未决率"]
ROWS = [
    ["Semgrep（同框架独立判定）", "57.4%", "42.3%", "57.5%", "—",     "70.0%", "0"],
    ["纯 LLM 零样本（qwen3-8b）", "96.7%", "26.9%", "89.7%", "84.8%", "37.5%¹", "0"],
    ["纯 LLM SFT（α0.5·combined）", "96.7%", "15.4%", "93.1%", "89.8%", "85.0%", "0"],
    ["两阶段融合（wave8）",       "100%", "4.3%", "98.7%", "92.3%", "94.1%", "13.8%"],
]
gf = s5.shapes.add_table(5, 7, Inches(0.75), Inches(1.72), Inches(7.9), Inches(2.35))
tbl = gf.table
tbl.first_row = False
tbl.horz_banding = False
widths = [2.30, 0.86, 0.86, 0.86, 1.02, 1.10, 0.90]
for i, w in enumerate(widths):
    tbl.columns[i].width = Inches(w)
tbl.rows[0].height = Inches(0.38)
for i in range(1, 5):
    tbl.rows[i].height = Inches(0.48)

def set_cell(cell, text, size, color, bold=False, fill=None, align=PP_ALIGN.CENTER, mono=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill if fill else WHITE
    cell.margin_left = Inches(0.05); cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.name = MONO if mono else FONT
    f.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT})
    rPr.append(ea)

# 表头
for j, h in enumerate(HEAD):
    set_cell(tbl.cell(0, j), h, 9.5, WHITE, bold=True, fill=BLUE_D)
# 数据行
HILITE = {3}
for i, row in enumerate(ROWS, start=1):
    hi = i in HILITE
    base_fill = BLUE_L if hi else (WHITE if i % 2 else MIST)
    for j, v in enumerate(row):
        if j == 0:
            set_cell(tbl.cell(i, j), v, 9.5, INK, bold=hi, fill=base_fill, align=PP_ALIGN.LEFT)
        else:
            col = INK
            bold = False
            if hi:
                bold = True
                if j in (2, 5): col = RED          # 误报率 / CVE-20 召回
                elif j == 4: col = BLUE_D           # strict
            elif j == 2 and i >= 2:
                col = RED                           # 误报率列淡红
            set_cell(tbl.cell(i, j), v, 9.5, col, bold=bold, fill=base_fill, mono=True)
print("对比表已插入")

# ---- 3) 表下要点 ----
def add_label(x, y, w, text, size=9.5, color=INK2, bold=False, name=FONT):
    tb = s5.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.name = name; f.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT}); rPr.append(ea)

add_label(0.78, 4.22, 7.8, "▎架构收益：误报 15.4%→4.3%（信任层收敛）｜strict 归因 89.8%→92.3%（证据链反哺）｜真实召回 85%→94.1%｜代价：13.8% 转人工", 10, BLUE_D, bold=True)
add_label(0.78, 4.52, 7.8, "▎SFT 收益：真实集召回 37.5%→85%（判别与格式）；但 strict 归因不随 SFT 提升，须由两阶段补足", 10, INK2)
add_label(0.78, 4.82, 7.8, "▎传统工具：召回 57.4% + 误报 42.3%，仅可作粗筛；典型样本双漏报（os.path.join 不做路径规范化）", 10, INK2)

# ---- 4) 右栏更新 ----
for sh in s5.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if t.startswith("EXP_02 · 14 段典型样本同题对打"):
        rewrite(sh, ["统一口径实测 · 87 段合成集（61 漏 / 26 安）"], 10, GOLD, bold=True)
    elif t == "Bandit":
        rewrite(sh, ["Semgrep"], 12.5, INK, bold=True)
    elif t.startswith("召回 83.3% 但误报率 50%"):
        rewrite(sh, ["召回 57.4% 但误报 42.3%：安全样本误杀近半，只能当粗筛"], 10.5, INK2)
    elif t == "Semgrep":
        rewrite(sh, ["纯 LLM（SFT 后）"], 12.5, INK, bold=True)
    elif t.startswith("误报为 0 但召回仅 75.0%"):
        rewrite(sh, ["零样本误报 26.9%，SFT 后降到 15.4%；但真实 CVE 召回不稳"], 10.5, INK2)
    elif t == "纯 LLM":
        rewrite(sh, ["两阶段融合（本工作）"], 12.5, INK, bold=True)
    elif t.startswith("14 段小样本满分 ≠ 真实可用"):
        rewrite(sh, ["误报 4.3% + strict 92.3% + 真实召回 94.1%，未决 13.8% 转人工"], 10.5, INK2)
    elif t.startswith("放大到 87 段合成集 + 20 段真实 CVE"):
        rewrite(sh, ["真实度阶梯（layer 6）"], 11.5, INK, bold=True)
    elif t.startswith("纯 LLM 合成集误报 15.4%"):
        rewrite(sh, [
            "合成集 87 → 真实 CVE 20 → rolling_dev 50：",
            "真实度上升、成绩下降（94.1% → 未正式跑），合成高召回不可外推",
        ], 10.5, INK2)
    elif t.startswith("结论：规则工具"):
        rewrite(sh, ["结论：误报 42%→15%→4.3%，真实召回 70%→85%→94%——分工架构被统一口径数据验证"], 12, BLUE_D, bold=True)
    elif t.startswith("口径：exp_02"):
        rewrite(sh, [
            "口径：acc=(TP+TN)/已裁决样本；strict=判真且 CWE 归因正确（2026-09-02 官方修正答案，score_batch 可复现，48 次运行全表见 unified_score_table_20260903）",
            "Semgrep 为 exp_07 融合评测同框架独立判定（stage1.by_tool 提取）；¹ 零样本 CVE-20 为 8 段小集；wave8 未决率 13.8%（12/87 转人工），全量口径 acc 85.1%",
            "两阶段引用注明组态：wave8=完整信任层（signal_feedback on + ctx16384）；CVE-fix 20 段全为真漏洞，无安全对照不计算 FPR",
        ], 8, INK3)

prs.save(SRC)
print("已保存")
