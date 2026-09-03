# -*- coding: utf-8 -*-
"""修复交付稿第5页（误报率不可见）与第6页（散点图混乱/版本错误）。
用法: python fix_p05_p06.py
数据口径：2026-09-03 联网核实（各厂官方发布 + swebench.com mini-SWE-agent v2.0.0 统一框架批次）。
"""
import copy, shutil, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SRC = "deliver/基于大语言模型的代码安全分析系统_30页修订版_20260903.pptx"
BAK = "deliver/backup_修复前_20260903.pptx"

# 调色板（与整套一致）
INK   = RGBColor(0x2B, 0x2F, 0x36)
INK2  = RGBColor(0x5B, 0x67, 0x70)
INK3  = RGBColor(0x94, 0x9E, 0xA6)
BLUE  = RGBColor(0x4A, 0x7F, 0xA5)
GOLD  = RGBColor(0xD9, 0x96, 0x2E)
RED   = RGBColor(0xC0, 0x50, 0x4D)
TEAL  = RGBColor(0x6F, 0xA3, 0x9B)
PAPER = RGBColor(0xFA, 0xF8, 0xF3)
MIST  = RGBColor(0xF4, 0xF7, 0xFA)
LINEC = RGBColor(0xE2, 0xE6, 0xEA)
BLUE_D= RGBColor(0x34, 0x5E, 0x80)
FONT  = "思源黑体"
MONO  = "Roboto Mono"

shutil.copyfile(SRC, BAK)
print("备份 ->", BAK)

prs = Presentation(SRC)

def set_font(run, size, color, bold=False, name=FONT):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.name = name
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', FONT)

def add_label(slide, x, y, w, text, size=7.5, color=INK2, name=MONO, align=PP_ALIGN.LEFT, bold=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.18))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    set_font(r, size, color, bold=bold, name=name)
    return tb

def add_dot(slide, cx, cy, d, fill=None, border=None, bw=1.5):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-d/2), Inches(cy-d/2), Inches(d), Inches(d))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.solid(); sp.fill.fore_color.rgb = PAPER
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if border is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = border; sp.line.width = Pt(bw)
    return sp

def add_hline(slide, x1, y1, x2, y2, color=LINEC, w=0.75, dash=None):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color; ln.line.width = Pt(w)
    ln.shadow.inherit = False
    if dash:
        lnEl = ln.line._get_or_add_ln()
        d = lnEl.makeelement(qn('a:prstDash'), {'val': dash}); lnEl.append(d)
    return ln

def rewrite_shape_text(sh, lines, size, color, bold=False, name=FONT, lh=None):
    """保留原首个 run 的字体样式，重写为多段文本。"""
    tf = sh.text_frame
    proto = None
    for p0 in tf.paragraphs:
        for r0 in p0.runs:
            proto = copy.deepcopy(r0._r)
            break
        if proto is not None:
            break
    # 清空所有段落，仅留第一段
    txBody = tf._txBody
    for p in txBody.findall(qn('a:p'))[1:]:
        txBody.remove(p)
    p1 = txBody.find(qn('a:p'))
    for child in list(p1):
        if child.tag != qn('a:pPr'):
            p1.remove(child)
    paras = [p1]
    for _ in range(len(lines) - 1):
        np = copy.deepcopy(p1)
        for child in list(np):
            if child.tag != qn('a:pPr'):
                np.remove(child)
        txBody.append(np); paras.append(np)
    for pel, line in zip(paras, lines):
        r = copy.deepcopy(proto)
        for t in r.findall(qn('a:t')):
            r.remove(t)
        import lxml.etree as etree
        t = etree.SubElement(r, qn('a:t')); t.text = line
        # 应用统一样式
        rPr = r.find(qn('a:rPr'))
        if rPr is None:
            rPr = r.makeelement(qn('a:rPr'), {'lang': 'zh-CN'}); r.insert(0, rPr)
        rPr.set('sz', str(int(size * 100)))
        rPr.set('b', '1' if bold else '0')
        for tag in ('a:solidFill',):
            for e in rPr.findall(qn(tag)):
                rPr.remove(e)
        import lxml.etree as etree2
        sf = etree2.SubElement(rPr, qn('a:solidFill'))
        clr = etree2.SubElement(sf, qn('a:srgbClr')); clr.set('val', str(color))
        pel.append(r)

# ================= 第 5 页：重建图表 =================
s5 = prs.slides[4]
old = None
for sh in list(s5.shapes):
    if sh.has_chart:
        old = sh
if old is not None:
    old._element.getparent().remove(old._element)
    print("第5页旧图表已移除")

cd = CategoryChartData()
cd.categories = ["纯 LLM 零样本", "Bandit", "Semgrep"]
cd.add_series("召回率", (1.0, 0.833, 0.75))
cd.add_series("误报率", (0.004, 0.5, 0.004))   # 0 值以极小柱体呈现，标签写 0.0%
cd.add_series("准确率", (1.0, 0.75, 0.786))

gf = s5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.75), Inches(1.67), Inches(7.5), Inches(4.55), cd)
ch = gf.chart
ch.has_title = False
ch.font.size = Pt(10); ch.font.name = FONT

# 去图表外框（spPr 按规范位于 c:chart 之后）
cs = ch._chartSpace
spPr = cs.makeelement(qn('c:spPr'), {})
ln = spPr.makeelement(qn('a:ln'), {})
ln.append(spPr.makeelement(qn('a:noFill'), {}))
spPr.append(ln)
chart_el = cs.find(qn('c:chart'))
chart_el.addnext(spPr)

colors = [BLUE, RED, TEAL]
plot = ch.plots[0]
plot.gap_width = 110
plot.overlap = -12
for i, ser in enumerate(plot.series):
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = colors[i]
    ser.format.line.fill.background()

plot.has_data_labels = True
dls = plot.data_labels
dls.show_value = True
dls.number_format = '0.0%'; dls.number_format_is_linked = False
dls.position = XL_LABEL_POSITION.OUTSIDE_END
dls.font.size = Pt(9); dls.font.color.rgb = INK2; dls.font.name = FONT

# 零值点自定义标签
zero_pts = [(1, 0), (1, 2)]  # (系列, 类目)：误报率的纯LLM、Semgrep
for si, ci in zero_pts:
    dl = plot.series[si].points[ci].data_label
    tf = dl.text_frame
    tf.text = "0.0%"
    r = tf.paragraphs[0].runs[0]
    set_font(r, 9, INK2)

va = ch.value_axis
va.minimum_scale = 0.0
va.maximum_scale = 1.08
va.major_unit = 0.2
va.tick_labels.number_format = '0%'
va.tick_labels.number_format_is_linked = False
va.tick_labels.font.size = Pt(9); va.tick_labels.font.color.rgb = INK3
va.has_major_gridlines = True
va.major_gridlines.format.line.color.rgb = LINEC
va.major_gridlines.format.line.width = Pt(0.75)
va.format.line.fill.background()

ca = ch.category_axis
ca.tick_labels.font.size = Pt(10.5)
ca.tick_labels.font.color.rgb = INK
ca.format.line.color.rgb = INK3
ca.has_major_gridlines = False

ch.has_legend = True
lg = ch.legend
lg.position = XL_LEGEND_POSITION.BOTTOM
lg.include_in_layout = False
lg.font.size = Pt(10); lg.font.color.rgb = INK2

# 图表下方补充说明（零值只显示标签）
add_label(s5, 0.75, 6.02, 7.5, "注：误报率为 0% 的方法无柱体，仅显示数值标签", size=8.5, color=INK3, name=FONT)
print("第5页新图表已插入")

# ================= 第 6 页：重绘左侧散点 =================
s6 = prs.slides[5]
kill = []
for sh in list(s6.shapes):
    L = Emu(sh.left).inches; T = Emu(sh.top).inches
    if L < 8.4 and 1.5 < T < 5.75:
        txt = sh.text_frame.text.strip() if sh.has_text_frame else ""
        if txt.startswith("SWE-bench Verified"):
            continue
        kill.append(sh)
for sh in kill:
    sh._element.getparent().remove(sh._element)
print(f"第6页移除旧散点元素 {len(kill)} 个")

# 几何
X0, Y0, PW, PH = 1.15, 1.90, 6.75, 2.85
VMIN, VMAX, MMAX = 25.0, 85.0, 26.0
def sx(m): return X0 + m / MMAX * PW
def sy(v): return Y0 + PH - (v - VMIN) / (VMAX - VMIN) * PH

# 绘图区底 + 网格 + 轴
bg = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(X0), Inches(Y0), Inches(PW), Inches(PH))
bg.shadow.inherit = False
bg.fill.solid(); bg.fill.fore_color.rgb = MIST
bg.line.fill.background()
for gv in range(30, 90, 10):
    gy = sy(gv)
    add_hline(s6, X0, gy, X0+PW, gy, color=LINEC, w=0.75)
    add_label(s6, X0-0.55, gy-0.075, 0.42, f"{gv}%", size=8, color=INK3, align=PP_ALIGN.RIGHT)
add_hline(s6, X0, Y0+PH, X0+PW, Y0+PH, color=INK3, w=1.0)
for m, lab in [(0,"2024.01"),(6,"2024.07"),(12,"2025.01"),(18,"2025.07"),(24,"2026.01")]:
    gx = sx(m)
    add_hline(s6, gx, Y0+PH, gx, Y0+PH+0.05, color=INK3, w=1.0)
    add_label(s6, gx-0.35, Y0+PH+0.07, 0.70, lab, size=8, color=INK3, align=PP_ALIGN.CENTER)

# 数据点：(标签, 分数, 月, 颜色, 空心, 标签位置模式)
# pos: ('R',dx,dy) 右侧 / ('L',dx,dy) 左侧(右对齐) / ('B',dx,dy) 下方
PTS = [
    ("GPT-4o 33.2",              33.2,  4,  BLUE, False, ('L', 0.08, -0.20)),
    ("Claude 3.5 Sonnet 33.4",   33.4,  5,  BLUE, False, ('B', 0.08,  0.04)),
    ("Claude 3.5 Sonnet(new) 49.0", 49.0, 9, BLUE, False, ('R', 0.08, -0.06)),
    ("DeepSeek-V3 42.0",         42.0, 11,  GOLD, False, ('B', 0.08,  0.04)),
    ("Gemini 2.5 Pro 63.8",      63.8, 14,  BLUE, False, ('B', 0.08,  0.04)),
    ("Claude Opus 4 72.5",       72.5, 16,  BLUE, False, ('L', 0.08, -0.20)),
    ("Qwen3-Coder 69.6",         69.6, 18,  GOLD, False, ('L', 0.08,  0.10)),
    ("Grok 4* 72.0",             72.0, 18,  BLUE, True,  ('R', 0.08, -0.08)),
    ("GPT-5 74.9",               74.9, 19,  BLUE, False, ('R', 0.08, -0.20)),
    ("GLM-4.6 68.0",             68.0, 21,  GOLD, False, ('R', 0.08,  0.08)),
    ("Claude Opus 4.5 80.9",     80.9, 22,  BLUE, False, ('L', 0.08, -0.06)),
]
for lab, v, m, col, hollow, mode in PTS:
    cx, cy = sx(m), sy(v)
    if hollow:
        add_dot(s6, cx, cy, 0.11, fill=PAPER, border=col, bw=1.5)
    else:
        add_dot(s6, cx, cy, 0.11, fill=col)
    side, dx, dy = mode
    w = 0.075 * len(lab)
    if side == 'R':
        add_label(s6, cx+dx, cy+dy, w, lab)
    elif side == 'L':
        add_label(s6, cx- dx - w, cy+dy, w, lab, align=PP_ALIGN.RIGHT)
    else:
        add_label(s6, cx+dx-0.0, cy+dy+0.10, w, lab)

# ---- 2026.02 统一框架批次（mini-SWE-agent v2.0.0，swebench.com）----
BATCH = [
    ("①", "Claude Opus 4.5 76.8", 76.8, BLUE, 7.50, 2.385),
    ("②", "Gemini 3 Flash 75.8",  75.8, BLUE, 7.58, 2.4325),
    ("③", "MiniMax M2.5 75.8",    75.8, GOLD, 7.64, 2.4325),
    ("④", "Claude Opus 4.6 75.6", 75.6, BLUE, 7.73, 2.4425),
    ("⑤", "GPT-5.2 72.8",         72.8, BLUE, 7.44, 2.558),
]
# 虚线分组框
box = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.36), Inches(2.31), Inches(0.47), Inches(0.33))
box.shadow.inherit = False
box.fill.background()
box.line.color.rgb = INK3; box.line.width = Pt(0.75)
lnEl = box.line._get_or_add_ln()
d = lnEl.makeelement(qn('a:prstDash'), {'val': 'dash'}); lnEl.append(d)

for no, name, v, col, cx, cy in BATCH:
    add_dot(s6, cx, cy, 0.10, fill=col)
# 编号（散放避让）
NUMPOS = {"①": (7.38, 2.44), "②": (7.525, 2.50), "③": (7.60, 2.62), "④": (7.72, 2.51), "⑤": (7.34, 2.62)}
for no, name, v, col, cx, cy in BATCH:
    nx, ny = NUMPOS[no]
    add_label(s6, nx, ny, 0.14, no, size=8, color=INK, name=FONT)

# 批次清单一行（图例上方）
add_label(s6, X0, 5.06, 7.6,
    "2026.02 统一框架批次（mini-SWE-agent，swebench.com）：① Claude Opus 4.5 76.8   ② Gemini 3 Flash 75.8   ③ MiniMax M2.5 75.8   ④ Claude Opus 4.6 75.6   ⑤ GPT-5.2 72.8",
    size=7.5, color=INK2, name=FONT)

# 图例
add_dot(s6, 1.20, 5.32, 0.09, fill=BLUE)
add_label(s6, 1.30, 5.27, 1.5, "海外模型", size=8.5, color=INK2, name=FONT)
add_dot(s6, 2.55, 5.32, 0.09, fill=GOLD)
add_label(s6, 2.65, 5.27, 1.5, "国产模型", size=8.5, color=INK2, name=FONT)
add_dot(s6, 3.80, 5.32, 0.09, fill=PAPER, border=INK3, bw=1.0)
add_label(s6, 3.90, 5.27, 2.0, "* 厂商自报框架", size=8.5, color=INK2, name=FONT)

# ---- 右侧结论与来源更新 ----
for sh in s6.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text
    if "国产进入第一梯队" in t:
        pass  # 标题不动
    elif "统一 harness 批次" in t and "GLM、Kimi" in t:
        rewrite_shape_text(sh,
            ["2026.02 统一框架批次五模型挤在 72.8\u201376.8；自报口径 MiniMax M2.5 80.2、GLM-5 77.8、Kimi K2.5 76.8，已进第一梯队"],
            10, INK2)
    elif t.startswith("来源：SWE-bench 官方榜单"):
        rewrite_shape_text(sh, [
            "来源：SWE-bench Verified 各厂官方发布（2026-09 检索）；统一框架为 mini-SWE-agent v2.0.0（swebench.com，2026-02-17），自报框架普遍高 4\u20136 分，标 * 分列",
            "官方分：GPT-4o 33.2 / Claude 3.5 Sonnet 33.4\u219249.0 / DeepSeek-V3 42.0 / Gemini 2.5 Pro 63.8 / Claude Opus 4 72.5 / Qwen3-Coder 69.6 / GPT-5 74.9 / GLM-4.6 68.0 / Claude Opus 4.5 80.9",
            "HumanEval 取各代官方报告（Codex 28.8 \u2192 Qwen2.5-Coder-32B 92.7）；Grok 4 为 xAI 自报口径（72\u201375）",
        ], 8.5, INK3)

# ---- HumanEval 图表修正（类目名 + 数据标签百分比）----
import zipfile, re, os
zin = zipfile.ZipFile(SRC)
target = None
for n in zin.namelist():
    if 'charts/chart' in n and n.endswith('.xml'):
        x = zin.read(n).decode('utf-8')
        if 'Codex' in x and 'HumanEval' in x:
            target = n; break
zin.close()
print("HumanEval 图表文件:", target)

prs.save(SRC)
print("已保存:", SRC)

# 后处理：直接改 zip 内 HumanEval chart xml
if target:
    tmp = SRC + ".tmp"
    zin = zipfile.ZipFile(SRC)
    zout = zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED)
    for item in zin.infolist():
        data = zin.read(item.filename)
        if item.filename == target:
            x = data.decode('utf-8')
            x = x.replace('<v>24 Qwen2.5</v>', '<v>24 Qwen2.5-Coder</v>')
            x = x.replace('formatCode="0.0"', 'formatCode="0.0%"')
            data = x.encode('utf-8')
        zout.writestr(item, data)
    zout.close(); zin.close()
    os.replace(tmp, SRC)
    print("HumanEval 图表类目与标签格式已修正")

print("ALL DONE")
