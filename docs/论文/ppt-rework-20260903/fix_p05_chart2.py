# -*- coding: utf-8 -*-
"""第5页：对比表 → 双柱状图（左：87段三指标×4方法；右：20段CVE召回）。
数据同 fix_p05_table.py（素材库官方口径 + exp_07 提取）。"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

SRC = "deliver/基于大语言模型的代码安全分析系统_30页修订版_20260903.pptx"

INK   = RGBColor(0x2B, 0x2F, 0x36)
INK2  = RGBColor(0x5B, 0x67, 0x70)
INK3  = RGBColor(0x94, 0x9E, 0xA6)
BLUE  = RGBColor(0x4A, 0x7F, 0xA5)
BLUE_D= RGBColor(0x34, 0x5E, 0x80)
GOLD  = RGBColor(0xD9, 0x96, 0x2E)
RED   = RGBColor(0xC0, 0x50, 0x4D)
TEAL  = RGBColor(0x6F, 0xA3, 0x9B)
GRAY  = RGBColor(0xB8, 0xC2, 0xCA)
LINEC = RGBColor(0xE2, 0xE6, 0xEA)
FONT  = "思源黑体"
MONO  = "Roboto Mono"

prs = Presentation(SRC)
s5 = prs.slides[4]

# ---- 1) 删除对比表 ----
for sh in list(s5.shapes):
    if sh.has_table:
        sh._element.getparent().remove(sh._element)
        print("对比表已删除")

def style_chart(ch, label_size=8, cat_size=8.5, legend=False):
    ch.has_title = False
    ch.font.size = Pt(label_size); ch.font.name = FONT
    cs = ch._chartSpace
    spPr = cs.makeelement(qn('c:spPr'), {})
    ln = spPr.makeelement(qn('a:ln'), {})
    ln.append(spPr.makeelement(qn('a:noFill'), {}))
    cs.insert(list(cs).index(cs.find(qn('c:chart'))) + 1, spPr)
    plot = ch.plots[0]
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.show_value = True
    dls.number_format = '0.0%'; dls.number_format_is_linked = False
    dls.position = XL_LABEL_POSITION.OUTSIDE_END
    dls.font.size = Pt(label_size); dls.font.color.rgb = INK2; dls.font.name = FONT
    va = ch.value_axis
    va.minimum_scale = 0.0; va.maximum_scale = 1.12; va.major_unit = 0.25
    va.tick_labels.number_format = '0%'
    va.tick_labels.number_format_is_linked = False
    va.tick_labels.font.size = Pt(8); va.tick_labels.font.color.rgb = INK3
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = LINEC
    va.major_gridlines.format.line.width = Pt(0.75)
    va.format.line.fill.background()
    ca = ch.category_axis
    ca.tick_labels.font.size = Pt(cat_size)
    ca.tick_labels.font.color.rgb = INK
    ca.format.line.color.rgb = INK3
    ca.has_major_gridlines = False
    ch.has_legend = legend
    if legend:
        lg = ch.legend
        lg.position = XL_LEGEND_POSITION.BOTTOM
        lg.include_in_layout = False
        lg.font.size = Pt(9); lg.font.color.rgb = INK2

# ---- 2) 左主图：87 段三指标 × 4 方法 ----
cd = CategoryChartData()
cd.categories = ["Semgrep", "纯LLM零样本", "纯LLM SFT", "两阶段融合"]
cd.add_series("召回率", (0.574, 0.967, 0.967, 1.0))
cd.add_series("误报率", (0.423, 0.269, 0.154, 0.043))
cd.add_series("准确率", (0.575, 0.897, 0.931, 0.987))
gf1 = s5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.75), Inches(1.70), Inches(4.90), Inches(2.62), cd)
ch1 = gf1.chart
style_chart(ch1, label_size=7.5, cat_size=8, legend=True)
colors = [BLUE, RED, TEAL]
plot1 = ch1.plots[0]
plot1.gap_width = 90
plot1.overlap = -10
for i, ser in enumerate(plot1.series):
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = colors[i]
    ser.format.line.fill.background()
print("左主图已插入")

# ---- 3) 右小图：20 段 CVE-fix 召回率 ----
cd2 = CategoryChartData()
cd2.categories = ["Semgrep", "纯LLM零样本", "纯LLM SFT", "两阶段融合"]
cd2.add_series("CVE-20 召回", (0.700, 0.375, 0.850, 0.941))
gf2 = s5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5.85), Inches(1.70), Inches(2.55), Inches(2.62), cd2)
ch2 = gf2.chart
style_chart(ch2, label_size=8, cat_size=7, legend=False)
plot2 = ch2.plots[0]
plot2.gap_width = 60
bar_colors = [GRAY, GRAY, BLUE, GOLD]
for i, pt in enumerate(plot2.series[0].points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = bar_colors[i]
    pt.format.line.fill.background()
print("右小图已插入")

# ---- 4) 图题 ----
def add_label(x, y, w, text, size=9.5, color=INK2, bold=True, name=FONT):
    tb = s5.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.name = name; f.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT}); rPr.append(ea)

add_label(0.78, 1.50, 4.9, "87 段合成集 · 召回 / 误报 / 准确（61 漏 / 26 安）", 9.5, BLUE_D)
add_label(5.88, 1.50, 2.6, "20 段真实 CVE-fix · 召回率", 9.5, BLUE_D)

# ---- 5) 右栏标题去重（原"同口径实测 · 87 段"改为方法解读）----
import copy
import lxml.etree as etree
def rewrite(sh, lines, size, color, bold=None):
    tf = sh.text_frame
    proto = None
    for p0 in tf.paragraphs:
        for r0 in p0.runs:
            proto = copy.deepcopy(r0._r); break
        if proto is not None: break
    txBody = tf._txBody
    for p in txBody.findall(qn('a:p'))[1:]: txBody.remove(p)
    p1 = txBody.find(qn('a:p'))
    for child in list(p1):
        if child.tag != qn('a:pPr'): p1.remove(child)
    paras = [p1]
    for _ in range(len(lines) - 1):
        np = copy.deepcopy(p1)
        for child in list(np):
            if child.tag != qn('a:pPr'): np.remove(child)
        txBody.append(np); paras.append(np)
    for pel, line in zip(paras, lines):
        r = copy.deepcopy(proto)
        for t in r.findall(qn('a:t')): r.remove(t)
        t = etree.SubElement(r, qn('a:t')); t.text = line
        rPr = r.find(qn('a:rPr'))
        if rPr is None:
            rPr = etree.Element(qn('a:rPr'), {'lang': 'zh-CN'}); r.insert(0, rPr)
        rPr.set('sz', str(int(size * 100)))
        if bold is not None: rPr.set('b', '1' if bold else '0')
        for e in rPr.findall(qn('a:solidFill')): rPr.remove(e)
        sf = etree.SubElement(rPr, qn('a:solidFill'))
        clr = etree.SubElement(sf, qn('a:srgbClr')); clr.set('val', str(color))
        pel.append(r)

GOLD = RGBColor(0xD9, 0x96, 0x2E)
for sh in s5.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if t.startswith("同口径实测 · 87 段"):
        rewrite(sh, ["方法解读（统一口径）"], 10, GOLD, bold=True)

prs.save(SRC)
print("已保存")
