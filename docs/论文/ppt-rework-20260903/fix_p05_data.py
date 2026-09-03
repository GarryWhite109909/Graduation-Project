# -*- coding: utf-8 -*-
"""第5页数据源更换：14段exp_02 → 87段合成集 + 20段真实CVE-fix（真实实验数据）。
数据出处：
- Semgrep 独立判定：exp_07 anchor 批次 wave8_ctx16384_final.20260902.json 的 stage1.by_tool 提取
- 纯 LLM 最优：exp_06_eval.finetuned_custom.combined_nosource.20260816_035517.json
- 融合架构：exp_07_full87.wave8_ctx16384_final.20260902.json（87段）/ cve_fix20.anchor_ctx16384.20260901.json（20段）
"""
import sys, io, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Pt
import lxml.etree as etree

SRC = "deliver/基于大语言模型的代码安全分析系统_30页修订版_20260903.pptx"
INK2 = RGBColor(0x5B, 0x67, 0x70)
INK3 = RGBColor(0x94, 0x9E, 0xA6)

def rewrite(sh, lines, size, color, bold=None):
    tf = sh.text_frame
    proto = None
    for p0 in tf.paragraphs:
        for r0 in p0.runs:
            proto = copy.deepcopy(r0._r); break
        if proto is not None: break
    txBody = tf._txBody
    for p in txBody.findall(qn('a:p'))[1:]: txBody.remove(p)
    p1 = txBody.find(qn('a:p'))
    for child in list(p1):
        if child.tag != qn('a:pPr'): p1.remove(child)
    paras = [p1]
    for _ in range(len(lines) - 1):
        np = copy.deepcopy(p1)
        for child in list(np):
            if child.tag != qn('a:pPr'): np.remove(child)
        txBody.append(np); paras.append(np)
    for pel, line in zip(paras, lines):
        r = copy.deepcopy(proto)
        for t in r.findall(qn('a:t')): r.remove(t)
        t = etree.SubElement(r, qn('a:t')); t.text = line
        rPr = r.find(qn('a:rPr'))
        if rPr is None:
            rPr = etree.Element(qn('a:rPr'), {'lang': 'zh-CN'}); r.insert(0, rPr)
        rPr.set('sz', str(int(size * 100)))
        if bold is not None: rPr.set('b', '1' if bold else '0')
        for e in rPr.findall(qn('a:solidFill')): rPr.remove(e)
        sf = etree.SubElement(rPr, qn('a:solidFill'))
        clr = etree.SubElement(sf, qn('a:srgbClr')); clr.set('val', str(color))
        pel.append(r)

prs = Presentation(SRC)
s5 = prs.slides[4]

# ---- 1) 图表换数据 ----
chart = None
for sh in s5.shapes:
    if sh.has_chart:
        chart = sh.chart
cd = CategoryChartData()
cd.categories = ["Semgrep", "纯 LLM 单阶段", "融合架构"]
cd.add_series("召回率", (0.574, 0.984, 1.0))
cd.add_series("误报率", (0.423, 0.231, 0.043))
cd.add_series("准确率", (0.575, 0.920, 0.987))
chart.replace_data(cd)
print("图表数据已更换")

# ---- 2) 右栏与底部文字 ----
for sh in list(s5.shapes):
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if t.startswith("EXP_02 · 14 段典型样本同题对打"):
        rewrite(sh, ["实测 · 87 段合成集（61 漏洞 / 26 安全）"], 10, RGBColor(0xD9, 0x96, 0x2E), bold=True)
    elif t == "Bandit":
        rewrite(sh, ["Semgrep"], 12.5, RGBColor(0x2B, 0x2F, 0x36), bold=True)
    elif t.startswith("召回 83.3% 但误报率 50%"):
        rewrite(sh, ["召回 57.4%，误报 42.3%：26 段安全样本误杀 11 段，规则泛化差"], 10.5, INK2)
    elif t == "Semgrep":
        rewrite(sh, ["纯 LLM 单阶段"], 12.5, RGBColor(0x2B, 0x2F, 0x36), bold=True)
    elif t.startswith("误报为 0 但召回仅 75.0%"):
        rewrite(sh, ["微调后单阶段最优：准确 92.0%，但误报仍 23.1%（误杀 6/26）"], 10.5, INK2)
    elif t == "纯 LLM":
        rewrite(sh, ["融合架构"], 12.5, RGBColor(0x2B, 0x2F, 0x36), bold=True)
    elif t.startswith("14 段小样本满分 ≠ 真实可用"):
        rewrite(sh, ["工具召回 + LLM 封闭裁决：误报压到 4.3%（1/23），准确 98.7%"], 10.5, INK2)
    elif t.startswith("放大到 87 段合成集 + 20 段真实 CVE"):
        rewrite(sh, ["20 段真实 CVE-fix 实测（全部为真漏洞）"], 11.5, RGBColor(0x2B, 0x2F, 0x36), bold=True)
    elif t.startswith("纯 LLM 合成集误报 15.4%"):
        rewrite(sh, [
            "<p>真实 CVE 召回：Semgrep 70% → 纯 LLM 85% → <span color=\"C0504D\" bold=\"true\" fontSize=\"11\">融合 94.1%</span></p>",
            "<p><span color=\"5B6770\" fontSize=\"10.5\">合成集上各家普遍虚高，真实集才是硬口径；融合确定判 16/17，3 段转人工复核</span></p>",
        ], 11, RGBColor(0x2B, 0x2F, 0x36))
    elif t.startswith("结论：规则工具"):
        rewrite(sh, ["结论：误报率 42.3% → 23.1% → 4.3%，召回 57.4% → 98.4% → 100%——工具找全、LLM 判准的分工被真实数据验证"], 12, RGBColor(0x34, 0x5E, 0x80), bold=True)
    elif t.startswith("口径：exp_02"):
        rewrite(sh, [
            "口径：87 段合成集（61 漏洞 / 26 安全）；Semgrep 为 exp_07 融合评测同框架独立判定（stage1.by_tool）；纯 LLM 取 exp_06 单阶段最优（combined_nosource，2026-08-16）",
            "融合架构取 exp_07 anchor 批次 wave8_ctx16384_final（2026-09-02）：recall/FPR/acc 为确定判定口径，12/87 转人工复核不计入分母；全量口径准确率 85.1%",
            "20 段 CVE-fix 全为真漏洞、无安全对照，不计算误报率；融合 16 TP / 1 FN / 3 review。数据可复现：score_batch / eval_two_stage",
        ], 8, INK3)
    elif t.startswith("注：误报率为 0%"):
        sh._element.getparent().remove(sh._element)
        print("已删除零值注释")

prs.save(SRC)
print("已保存")
